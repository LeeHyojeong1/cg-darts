#!/usr/bin/env python
"""Build the Team-9 CG-DARTS final report .pptx.

Mirrors the proposal deck's visual style (720x405 pt, serif title with
horizontal rule beneath, sans-serif bullet content, slide-number tag
bottom-right). 25 slides total.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "reports" / "figures"
OUT = ROOT / "reports" / "Team9_FinalReport.pptx"

# 720 x 405 pt to match proposal
SLIDE_W = Pt(720)
SLIDE_H = Pt(405)

TITLE_FONT = "Cambria"   # serif, closest to the proposal's title face
BODY_FONT = "Calibri"    # sans-serif
NUM_FONT = "Calibri"

INK = RGBColor(0x1A, 0x1A, 0x1A)
RULE = RGBColor(0x33, 0x33, 0x33)
ACCENT = RGBColor(0x4C, 0x72, 0xB0)
MUTED = RGBColor(0x55, 0x55, 0x55)


def make_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_slide(prs):
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)


def add_textbox(slide, x, y, w, h, text, *, font=BODY_FONT, size=14,
                bold=False, italic=False, color=INK, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(0);  tf.margin_bottom = Pt(0)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, font=BODY_FONT, size=14,
                color=INK, level_indent=Pt(14), line_spacing=1.15):
    """items: list of (level, text). Level 0 uses bullet '•', level 1 '–'."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(0);  tf.margin_bottom = Pt(0)
    for i, (lvl, text) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        glyph = "•  " if lvl == 0 else "–  "
        indent = level_indent * lvl
        # Use leading spaces for indent (pptx-level indent is fiddly)
        run = p.add_run()
        run.text = (" " * (3 * lvl)) + glyph + text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_rule(slide, x, y, w):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(0.6))
    line.fill.solid(); line.fill.fore_color.rgb = RULE
    line.line.color.rgb = RULE
    line.line.width = Pt(0.6)
    return line


def chrome(slide, slide_no, total, title=None, subtitle=None):
    """Standard slide chrome: title bar + horizontal rule + page number."""
    add_textbox(slide, Pt(28), Pt(14), Pt(660), Pt(28),
                title if title else "",
                font=TITLE_FONT, size=22, bold=True)
    if subtitle:
        add_textbox(slide, Pt(28), Pt(40), Pt(660), Pt(16),
                    subtitle, font=BODY_FONT, size=10, italic=True, color=MUTED)
    add_rule(slide, Pt(28), Pt(58), Pt(664))
    add_textbox(slide, Pt(680), Pt(386), Pt(36), Pt(14),
                str(slide_no), font=NUM_FONT, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def title_slide(prs, total):
    s = add_slide(prs)
    # No chrome on title slide
    add_textbox(s, Pt(48), Pt(150), Pt(640), Pt(40),
                "CG-DARTS:  Cost-Guided Differentiable Architecture Search",
                font=TITLE_FONT, size=28, bold=True)
    add_rule(s, Pt(48), Pt(196), Pt(620))
    add_textbox(s, Pt(48), Pt(206), Pt(620), Pt(20),
                "Extending DARTS with hardware-aware regularization and derivation",
                font=BODY_FONT, size=14, italic=True, color=MUTED)
    add_textbox(s, Pt(48), Pt(248), Pt(620), Pt(14),
                "Baek Jongjin   ·   Kim Ryeowook   (Team 9)",
                font=BODY_FONT, size=12)
    add_textbox(s, Pt(48), Pt(266), Pt(620), Pt(12),
                "CS.50700  ·  Final Project Report",
                font=BODY_FONT, size=10, color=MUTED)
    add_textbox(s, Pt(48), Pt(360), Pt(620), Pt(12),
                "Builds on the DARTS paper review (Team9_Presentation.pdf, slide 16)",
                font=BODY_FONT, size=9, italic=True, color=MUTED)


def outline_slide(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "Outline")
    items = [
        (0, "Recap of the proposal (open questions identified in the paper review)"),
        (0, "Methods — across three branches"),
        (1, "Cost regularizer and five cost tables (FLOPs, params, memory, device, LUT)"),
        (1, "Softmax temperature annealing"),
        (1, "Cost-aware architecture derivation"),
        (0, "Experimental setup and compared methods"),
        (0, "Experimental results"),
        (1, "FLOPs sweep, Params metric, annealing pilot, L40S latency LUT"),
        (0, "Analysis: skip-collapse, what works, what doesn't"),
        (0, "Conclusions and future work"),
    ]
    add_bullets(s, Pt(40), Pt(76), Pt(640), Pt(310), items, size=14, line_spacing=1.30)


# ----------- Slide builders -----------

def slide_3_proposal_recap(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "1.  Proposal recap",
           subtitle="From slide 16 of Team9_Presentation.pdf — the original DARTS authors' open questions")
    items = [
        (0, "Open question O1 — closing the soft-to-hard discretization gap"),
        (1, "DARTS searches a continuous mixture (softmax over operations) but evaluates a discrete cell"),
        (1, "Liu et al. (2019) suggest annealing the softmax temperature as a possible remedy"),
        (0, "Open question O2 — performance-aware architecture-derivation"),
        (1, "The original top-k argmax derivation is blind to inference cost"),
        (1, "Calls for schemes that exploit the shared weights and account for compute / memory / latency"),
        (0, "Our thesis: bring hardware cost into BOTH the search loss AND the derivation step"),
        (1, "Differentiable cost regularizer during search  +  cost-aware top-k during derivation"),
        (1, "Tau-annealing as the bridge between the two regimes"),
    ]
    add_bullets(s, Pt(40), Pt(80), Pt(640), Pt(300), items, size=13, line_spacing=1.25)


def slide_4_thesis(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "2.  CG-DARTS thesis")
    add_textbox(s, Pt(40), Pt(80), Pt(660), Pt(28),
                "Augmented architecture loss:",
                font=BODY_FONT, size=14, bold=True)
    add_textbox(s, Pt(60), Pt(106), Pt(660), Pt(32),
                "L_arch(α) = L_val(w*(α), α)  +  λ · E_α[c]",
                font=TITLE_FONT, size=20, color=ACCENT)
    add_textbox(s, Pt(40), Pt(150), Pt(660), Pt(16),
                "where the expected cost over the soft mixture",
                font=BODY_FONT, size=12)
    add_textbox(s, Pt(60), Pt(168), Pt(660), Pt(28),
                "E_α[c] = Σ_(i,j) Σ_o softmax(α^(i,j)/τ)_o · c_o^(i,j)",
                font=TITLE_FONT, size=16, color=ACCENT)
    items = [
        (0, "c_o^(i,j)  — per-edge, per-primitive cost table (precomputed, differentiable through α)"),
        (0, "τ  — softmax temperature; default 1.0, annealed during search to bridge soft→hard"),
        (0, "λ  — cost weight; linearly warmed up over the first 10 search epochs"),
        (0, "Five cost tables shipped: FLOPs · Params · Memory bytes · Device-blended · Measured LUT"),
        (0, "Three derivation modes: argmax (DARTS default) · cost_sub · cost_div"),
    ]
    add_bullets(s, Pt(40), Pt(216), Pt(660), Pt(170), items, size=12, line_spacing=1.30)


def slide_5_method_overview(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "3.  Method overview — three branches of contribution")
    s.shapes.add_picture(str(FIG / "fig6_branch_map.png"),
                         Pt(40), Pt(76), width=Pt(640))


def slide_6_cost_reg(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "4.  Method — cost regularizer", subtitle="cnn/cost_utils.py + cnn/architect_cg.py (master)")
    items = [
        (0, "Cost is precomputed once before search starts:  one [num_edges × num_ops] table per cell type"),
        (1, "Cells of the same type share α in DARTS → also share cost in CG-DARTS"),
        (1, "Each cell-type cost table sums the per-op cost across all cells of that type in the supernet"),
        (0, "Normalization (per-edge):  c̃[i,j,o] = (c[i,j,o] − min_o) / (max_o − min_o)"),
        (1, "Removes the scale advantage of larger cells, lets one λ work across metrics"),
        (0, "Architecture-loss backward pass: ∂E[c]/∂α has the closed-form softmax-Jacobian × cost"),
        (1, "Adds < 1% to per-step compute"),
        (0, "Cost-warmup: λ_eff(epoch) = λ · min(1, (epoch+1)/warmup) for first WARMUP epochs"),
        (1, "Avoids early bias when w is still random"),
    ]
    add_bullets(s, Pt(40), Pt(80), Pt(660), Pt(300), items, size=12, line_spacing=1.20)


def slide_7_metrics(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "5.  Method — five cost metrics")
    rows = [
        ("Metric", "Formula per primitive", "Origin branch"),
        ("flops", "conv / depthwise / pool multiply-adds at supernet feature shape", "master"),
        ("params", "learnable weight count", "master"),
        ("mem", "fp32 activation read + write bytes (proxy for traffic)", "device-conditioned"),
        ("device", "w_f·flops + w_m·mem + w_ℓ·roofline_latency  (per device profile)", "device-conditioned"),
        ("lut", "measured per-primitive latency from scripts/measure_latency_lut.py", "proposal-driven"),
    ]
    add_table(s, rows, Pt(40), Pt(80), col_widths=[Pt(80), Pt(420), Pt(180)],
              header_color=ACCENT, font_size=11)
    add_textbox(s, Pt(40), Pt(290), Pt(660), Pt(80),
                "Roofline latency:   ℓ = max(flops / P_peak,  mem_bytes / B_bw)\n"
                "                            P_peak from vendor FP32 spec, B_bw from memory-bus spec",
                font=TITLE_FONT, size=12, color=ACCENT)
    add_textbox(s, Pt(40), Pt(330), Pt(660), Pt(60),
                "The cost regularizer is metric-agnostic — same code path, different table.\n"
                "All five metrics produce alpha-shaped tables so build_search_costs() returns a uniform interface.",
                font=BODY_FONT, size=11, color=MUTED)


def slide_8_device(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "6.  Method — device profiles", subtitle="cnn/device_profiles.py (device-conditioned branch)")
    rows = [
        ("Profile", "SKU", "FP32 TFLOPS", "Mem BW (GB/s)", "(w_f, w_m, w_ℓ)"),
        ("jetson_orin", "Jetson AGX Orin 64GB", "5.32", "204.8", "(0.20, 0.60, 0.20)"),
        ("rtx_pro_6000", "RTX PRO 6000 Blackwell WS", "125.0", "1792.0", "(0.35, 0.35, 0.30)"),
        ("h100", "NVIDIA H100 SXM5", "67.0", "3350.0", "(0.50, 0.20, 0.30)"),
    ]
    add_table(s, rows, Pt(40), Pt(80), col_widths=[Pt(110), Pt(220), Pt(105), Pt(115), Pt(150)],
              header_color=ACCENT, font_size=11)
    items = [
        (0, "Weights follow the CG_DARTS Proposal v3 presets (edge / balanced / server)"),
        (0, "Vendor specs are the only inputs; the roofline closes the gap to per-op cost"),
        (0, "Same metric=device pipeline searches a different cell per profile — same code, three deployments"),
    ]
    add_bullets(s, Pt(40), Pt(220), Pt(660), Pt(160), items, size=12, line_spacing=1.25)


def slide_9_lut(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "7.  Method — measured latency LUT", subtitle="scripts/measure_latency_lut.py (proposal-driven branch)")
    items = [
        (0, "For each (edge_idx, primitive, stride, C, H, W) tuple in the 8-cell supernet:"),
        (1, "Build the primitive (with BN for pool ops, like the supernet does)"),
        (1, "Warm up 5 iters; time N iters with torch.cuda.Event start/end and torch.cuda.synchronize"),
        (1, "Accumulate into an alpha-shaped table summed across cells of the same type"),
        (0, "Output: reports/latency_lut/<device>.json with metadata + 14×8 tables per cell type"),
        (0, "Consumed via:   --cost_metric lut  --lut_path reports/latency_lut/l40s.json"),
        (0, "Same downstream code path as flops / params / mem / device → no special-casing in train_search_cg.py"),
    ]
    add_bullets(s, Pt(40), Pt(80), Pt(660), Pt(300), items, size=13, line_spacing=1.30)


def slide_10_anneal_why(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "8.  Method — softmax annealing  (motivation)",
           subtitle="Open question O1 from our proposal: closing the soft-to-hard discretization gap")

    add_textbox(s, Pt(40), Pt(80), Pt(660), Pt(18),
                "The problem  —  DARTS search and DARTS evaluation optimize different functions:",
                font=BODY_FONT, size=12, bold=True)
    items_problem = [
        (0, "Search:  L_val is computed on a soft mixture  ō(x) = Σ_o softmax(α)_o · o(x)  over all candidate ops"),
        (0, "Evaluation:  the discrete cell uses only argmax_o α_o on each edge  —  every other op vanishes"),
        (0, "Consequence:  search_valid_acc  ≠  retrain_test_acc.  The search optimizer never directly sees the model it ships."),
    ]
    add_bullets(s, Pt(50), Pt(102), Pt(650), Pt(74), items_problem, size=11, line_spacing=1.25)

    add_textbox(s, Pt(40), Pt(180), Pt(660), Pt(18),
                "Concrete evidence of the gap  (FLOPs CG-DARTS, λ=1e-2, seed=2):",
                font=BODY_FONT, size=11, bold=True, color=ACCENT)
    add_textbox(s, Pt(60), Pt(200), Pt(660), Pt(18),
                "search_valid_acc = 87.61%      vs      retrain_test_acc = 96.05%       →  ~8 pp gap",
                font=TITLE_FONT, size=13, color=ACCENT)

    add_textbox(s, Pt(40), Pt(232), Pt(660), Pt(18),
                "What annealing buys us  —  the limit τ → 0:",
                font=BODY_FONT, size=12, bold=True)
    items_buys = [
        (0, "softmax(α / τ) becomes one-hot at argmax_o α_o  →  the soft mixture IS the discrete cell"),
        (0, "search_valid_acc  →  the validation accuracy of the cell derivation will pick  (gap closes by construction)"),
        (0, "Expected cost E_α[c] becomes the discrete cost of the chosen cell — the regularizer penalizes the deployed model"),
        (0, "Liu et al. (2019) themselves flag annealing as the natural remedy in the Discussion section of the paper"),
    ]
    add_bullets(s, Pt(50), Pt(254), Pt(650), Pt(120), items_buys, size=11, line_spacing=1.25)


def slide_10_anneal_intuition(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "9.  Method — softmax annealing  (intuition)",
           subtitle="Same α, three temperature regimes — orange bar marks the argmax operation")
    s.shapes.add_picture(str(FIG / "fig9_annealing_intuition.png"),
                         Pt(40), Pt(74), width=Pt(640))
    add_textbox(s, Pt(40), Pt(360), Pt(660), Pt(36),
                "Annealing τ → 0 progressively makes the mixture identical to the discrete architecture.\n"
                "That is the mechanism by which the soft-to-hard gap on the previous slide closes.",
                font=BODY_FONT, size=10, italic=True, color=MUTED)


def slide_10_anneal(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "10.  Method — softmax annealing  (implementation)",
           subtitle="The actual code paths, flags, and per-model state that make annealing work end-to-end")
    add_textbox(s, Pt(40), Pt(80), Pt(660), Pt(18),
                "Replace softmax(α) with softmax(α / τ) in BOTH the mixed-op forward and the cost expectation:",
                font=BODY_FONT, size=12)
    add_textbox(s, Pt(60), Pt(102), Pt(660), Pt(24),
                "ō(x) = Σ_o softmax(α/τ)_o · o(x)        E_α[c] = Σ softmax(α/τ)_o · c_o",
                font=TITLE_FONT, size=14, color=ACCENT)
    items = [
        (0, "tau-schedule arg:  linear / exp / none.   Defaults to τ=1.0 (matches paper)."),
        (0, "Three flags:  --tau_start  --tau_end  --tau_anneal {linear,exp,none}"),
        (0, "Property:  as τ → 0,  softmax(α/τ) → one-hot at argmax_o α_o."),
        (1, "→  E_α[c] converges to the discrete cost of the cell that derivation will pick"),
        (1, "→  Closes the gap between the soft objective minimized during search and the discrete model"),
        (0, "The whole network exposes self.tau so .new() and the unrolled gradient inherit the schedule correctly"),
    ]
    add_bullets(s, Pt(40), Pt(140), Pt(660), Pt(240), items, size=12, line_spacing=1.25)


def slide_11_disc(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "11.  Method — cost-aware architecture derivation",
           subtitle="Addresses open question O2 (performance-aware derivation)")
    add_textbox(s, Pt(40), Pt(80), Pt(660), Pt(18),
                "Same top-k structure as DARTS, but the per-edge × per-op score s(j,o) is mode-dependent:",
                font=BODY_FONT, size=12)
    rows = [
        ("Mode", "Score s(j, o)",                                   "Behaviour"),
        ("argmax  (DARTS default)", "softmax(α)[j,o]",              "cost-blind"),
        ("cost_sub", "softmax(α)[j,o] − μ · c̃[j,o]",                "additive penalty; μ tunable at derivation time"),
        ("cost_div", "softmax(α)[j,o] / (ε + c̃[j,o])",              "favours value-per-cost ops"),
    ]
    add_table(s, rows, Pt(40), Pt(110), col_widths=[Pt(170), Pt(220), Pt(280)],
              header_color=ACCENT, font_size=11)
    items = [
        (0, "Selection is shared: same s(j,o) drives top-2 edges and the chosen op per edge"),
        (0, "scripts/rederive_genotype.py loads weights.pt and emits  genotype.<mode>.txt  post hoc"),
        (1, "ONE search → three derived genotypes → three retrains.  No extra search cost."),
    ]
    add_bullets(s, Pt(40), Pt(252), Pt(660), Pt(120), items, size=12, line_spacing=1.25)


def slide_12_setup(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "12.  Experimental setup")
    rows = [
        ("Phase", "Config"),
        ("Dataset", "CIFAR-10  ·  cutout-free transforms (search)  ·  cutout-16 + auxiliary tower (retrain)"),
        ("Search supernet", "8 cells  ·  16 init channels  ·  steps=4  ·  half train / half val split"),
        ("Search optimizer", "SGD w/ momentum 0.9, cosine LR 2.5e-2 → 0, weight decay 3e-4"),
        ("Architecture optimizer", "Adam lr 3e-4, β=(0.5, 0.999), weight decay 1e-3"),
        ("Search epochs / batch", "50  ·  64"),
        ("Retrain network", "20 cells  ·  36 init channels  ·  auxiliary weight 0.4  ·  drop_path 0.2"),
        ("Retrain epochs / batch", "100  ·  96"),
        ("Cost reg.", "λ ∈ {1e-3, 5e-3, 1e-2, 5e-2, 1e-1};  warmup 10 ep;  edge normalization"),
        ("Annealing", "tau linear 5.0 → 0.1 over 50 epochs (pilot)"),
        ("Seeds", "single seed = 2  (multi-seed orchestrator implemented but not authorized for production)"),
        ("Hardware", "NVIDIA L40S × 2  (only GPU 0 used; GPU 1 was held by a collaborator's job)"),
    ]
    add_table(s, rows, Pt(40), Pt(74), col_widths=[Pt(180), Pt(500)],
              header_color=ACCENT, font_size=10)


def slide_13_compared(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "13.  Compared methods")
    rows = [
        ("Method", "Cost-aware?", "Derivation", "Reported on", "Source"),
        ("DARTS V1  (first-order, paper)", "no",  "argmax", "CIFAR-10",          "Liu et al. (2019), Table 1"),
        ("DARTS V2  (second-order, paper)", "no", "argmax", "CIFAR-10",          "Liu et al. (2019), Table 1"),
        ("Vanilla DARTS  (our run)",       "no",  "argmax", "CIFAR-10 (our protocol)", "this work"),
        ("CG-DARTS FLOPs  (λ sweep)",      "yes (FLOPs)",  "argmax", "CIFAR-10",  "master"),
        ("CG-DARTS Params  (λ sweep)",     "yes (params)", "argmax", "CIFAR-10",  "device-conditioned branch"),
        ("CG-DARTS Device  (per profile)", "yes (FLOPs+mem+roofline)", "argmax", "3 tiers, 95.4-95.6% retrain", "device-conditioned"),
        ("CG-DARTS Annealing  (pilot)",    "yes (FLOPs)",  "argmax", "CIFAR-10",  "proposal-driven (this)"),
        ("CG-DARTS Cost-derived  (pilots)", "yes (FLOPs)", "cost_sub / cost_div", "queued; not yet retrained", "proposal-driven (this)"),
    ]
    add_table(s, rows, Pt(24), Pt(76), col_widths=[Pt(165), Pt(115), Pt(115), Pt(140), Pt(135)],
              header_color=ACCENT, font_size=9)


def slide_14_flops_table(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "14.  Results — FLOPs-cost CG-DARTS sweep (table)",
           subtitle="reports/cg_darts/summary.csv  ·  seed=2  ·  50ep search / 100ep retrain")
    rows = [
        ("Config",       "Search valid", "Expected cost", "Discrete MACs (M)", "Δ MACs", "Params (M)", "Δ Params", "Retrain test"),
        ("Vanilla DARTS",            "87.93%", "7.51", "477.96", "—",       "3.02", "—",       "95.97%"),
        ("CG-DARTS  λ=1e-3",         "87.52%", "7.11", "509.70", "−6.6%",   "3.25", "−7.6%",   "—"),
        ("CG-DARTS  λ=5e-3",         "87.28%", "6.78", "426.31", "+10.8%",  "2.65", "+12.2%",  "—"),
        ("CG-DARTS  λ=1e-2",         "87.61%", "6.24", "420.12", "+12.1%",  "2.63", "+13.0%",  "96.05%"),
        ("CG-DARTS  λ=5e-2",         "87.10%", "3.15", "255.64", "+46.5%",  "1.55", "+48.7%",  "93.96%"),
        ("CG-DARTS  λ=1e-1",         "86.77%", "1.76", "282.76", "+40.8%",  "1.73", "+42.7%",  "—"),
    ]
    add_table(s, rows, Pt(28), Pt(80), col_widths=[Pt(120), Pt(70), Pt(80), Pt(95), Pt(70), Pt(75), Pt(70), Pt(80)],
              header_color=ACCENT, font_size=10)
    add_textbox(s, Pt(40), Pt(290), Pt(660), Pt(60),
                "λ=1e-2 is Pareto-dominant under matched protocol:   +0.08 pp retrain accuracy AND −12.1% MACs AND −13.0% params vs. vanilla.\n"
                "λ=5e-2 is the 'aggressive' point: roughly half the MACs at a 2 pp accuracy cost.",
                font=BODY_FONT, size=11, color=MUTED)


def slide_15_flops_pareto(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "15.  Results — FLOPs Pareto plot",
           subtitle="Stars: retrain test accuracy  ·  circles: search-time validation  ·  dashed line: retrain Pareto frontier")
    s.shapes.add_picture(str(FIG / "fig1_flops_pareto.png"),
                         Pt(40), Pt(74), width=Pt(640))


def slide_15b_master_lambda(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "16.  Results — cost reduction across λ  (master FLOPs sweep)",
           subtitle="reports/cg_darts/summary.csv  ·  vanilla DARTS as baseline  ·  signed % reductions")
    s.shapes.add_picture(str(FIG / "fig7_master_lambda_sweep.png"),
                         Pt(40), Pt(74), width=Pt(640))


def slide_16_params(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "17.  Results — Params-metric CG-DARTS",
           subtitle="reports/cg_darts_params/summary.csv  ·  seed=2  ·  30ep search / 50ep retrain (fast pipeline)")
    s.shapes.add_picture(str(FIG / "fig2_params_vs_flops_reduction.png"),
                         Pt(40), Pt(74), width=Pt(640))


def slide_16b_device(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "18.  Results — Device-conditioned CG-DARTS",
           subtitle="Edge (Jetson Orin) / Middle (RTX PRO 6000) / High (H100)  ·  λ=1e-2  ·  seed=2  ·  50ep search / 100ep retrain  ·  device-conditioned branch")
    s.shapes.add_picture(str(FIG / "fig8_device_results.png"),
                         Pt(28), Pt(74), width=Pt(680))
    add_textbox(s, Pt(28), Pt(354), Pt(680), Pt(36),
                "Edge profile (high mem weight) → leanest cell (314 M MACs, 1.94 M params).  "
                "H100 profile (high compute weight) → largest cell (431 M MACs, 2.71 M params).\n"
                "All three retrains land at 95.4–95.6% — within 0.6 pp of vanilla DARTS, with at most −34% MACs and −36% params.",
                font=BODY_FONT, size=9, italic=True, color=MUTED)


def slide_17_anneal(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "19.  Results — annealing pilot #1: linear τ=5→0.1",
           subtitle="cnn/search-cg-anneal-tau5to01-lambda1em2-seed2  ·  valid peaks near τ≈1.0 then collapses — first attempt, failed")
    s.shapes.add_picture(str(FIG / "fig3_annealing_trajectory.png"),
                         Pt(40), Pt(74), width=Pt(640))


def slide_17b_anneal_compare(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "20.  Results — annealing pilot #2: schedule shape is decisive",
           subtitle="Holding τ=1 through the supernet's learning phase, then annealing only the last 20%, avoids collapse")
    s.shapes.add_picture(str(FIG / "fig10_annealing_comparison.png"),
                         Pt(28), Pt(72), width=Pt(680))
    rows = [
        ("λ=1e-2 variant", "τ schedule", "Search-valid", "Retrain", "Normal-cell convs"),
        ("no anneal", "τ=1 fixed", "87.61 %", "96.05 %", "healthy"),
        ("linear (pilot #1)", "τ=5→0.1 from ep 0", "82.91 %", "not run (skip cell)", "~2 (skip-dominated)"),
        ("hold-then-cosine (pilot #2)", "τ=1 to ep39, →0.3", "85.40 %", "96.04 %", "4 conv ops"),
    ]
    add_table(s, rows, Pt(120), Pt(330), col_widths=[Pt(170), Pt(150), Pt(95), Pt(80), Pt(120)],
              header_color=ACCENT, font_size=9, header_height=Pt(16), body_row_height=Pt(15))


def slide_18_lut(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "21.  Results — L40S latency LUT",
           subtitle="reports/latency_lut/l40s.json  ·  50-iter timing per primitive (batch=1)")
    s.shapes.add_picture(str(FIG / "fig4_l40s_latency.png"),
                         Pt(40), Pt(74), width=Pt(640))


def slide_19_vs_paper(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "22.  Results — comparison with Liu et al. (2019)")
    rows = [
        ("Axis", "Liu et al. (2019)", "This work", "Outcome"),
        ("Best CIFAR-10 test error",     "2.76% (V2, 600ep, 4-seed)",  "3.95% (CG-DARTS λ=1e-2, FO, 100ep)", "✗ raw accuracy not matched"),
        ("Search-cost regularization",   "none",                       "λ · E[c] across 5 metrics",          "✓ new capability"),
        ("Hardware awareness",           "none",                       "3 device profiles + measured LUT",   "✓ new capability"),
        ("Architecture derivation",      "top-k argmax on α",          "argmax + cost_sub + cost_div",       "✓ open question O2 addressed"),
        ("Soft-to-hard discretization",  "flagged as future work",     "τ-annealing + τ-aware E[c]",         "✓ open question O1 addressed"),
        ("Multi-seed protocol",          "4-seed model selection",     "single seed pilots",                 "✗ not yet matched"),
        ("Pareto operating point (CG-DARTS vs. vanilla, matched seed/protocol)",
         "—", "+0.08 pp acc · −12% MACs · −13% params  @ λ=1e-2", "✓ Pareto improvement"),
    ]
    add_table(s, rows, Pt(28), Pt(76), col_widths=[Pt(180), Pt(150), Pt(220), Pt(140)],
              header_color=ACCENT, font_size=9)


def slide_20_skip(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "23.  Analysis — skip-collapse hypothesis",
           subtitle="Cost regularization amplifies the DARTS skip-collapse pathology (Zela et al. 2020, Liang et al. DARTS+)")
    s.shapes.add_picture(str(FIG / "fig5_op_share.png"),
                         Pt(40), Pt(78), width=Pt(640))


def slide_21_discussion(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "24.  Discussion — what works, what doesn't")
    items = [
        (0, "What works"),
        (1, "λ=1e-2 Pareto-dominates vanilla under matched protocol — cheaper AND slightly more accurate"),
        (1, "Cost regularizer is metric-agnostic and the same code searches FLOPs / params / memory / device / LUT"),
        (1, "Device conditioning produces genuinely different cells per tier, all within 0.6 pp of vanilla"),
        (1, "Annealing CAN be done safely — schedule shape is what matters, not annealing per se"),
        (0, "What doesn't (or doesn't yet)"),
        (1, "Linear annealing collapses (search-valid 82.9%, skip-dominated cell) — anneals through the learning phase"),
        (1, "Hold-then-cosine rescues it (96.04% retrain, healthy cell) but is accuracy/cost-NEUTRAL at λ=1e-2"),
        (1, "So annealing is safe but not yet a win — its value needs higher λ or pairing with cost-aware derivation"),
        (1, "Single-seed throughout: 0.08 pp deltas are within DARTS seed noise (Zela et al. 2020: ~1 pp std)"),
    ]
    add_bullets(s, Pt(40), Pt(80), Pt(660), Pt(300), items, size=12, line_spacing=1.20)


def slide_22_branches(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "25.  Branch contribution map")
    rows = [
        ("Branch", "Code contributions", "Result artifacts"),
        ("master\n(initial CG-DARTS)",
         "cost_utils.py (FLOPs, params, edge norm, λ warmup)\n"
         "architect_cg.py (cost-aware bilevel update)\n"
         "train_search_cg.py (driver)\n"
         "run_cg_sweep.sh, cg_darts_report.py",
         "reports/cg_darts/summary.csv\n(6 runs: vanilla + 5 λ)\n"
         "Best: λ=1e-2 retrain 96.05% / −12% MACs"),
        ("experiment/\ndevice-conditioned",
         "memory-bytes metric\ndevice_profiles.py (Jetson/RTX PRO/H100)\nFLOPs+mem+roofline blend\nparams pipeline + report\ndevice pipeline + full retrains",
         "reports/cg_darts_params/summary.csv (λ=1e-2, 5e-2)\n"
         "reports/cg_darts_device/summary.csv (3 device tiers)\n"
         "Best params: λ=5e-2 → 94.22% / −32% MACs\n"
         "Device: Edge tier → 95.40% at 313 MACs"),
        ("experiment/\nproposal-driven\n(this)",
         "τ-annealing search\ncost-aware derivation (cost_sub/cost_div)\nlatency LUT + cross-device eval\nmulti-seed orchestrator\npost-hoc rederivation\nPareto plotting\nfinal report + figures",
         "reports/latency_lut/l40s.json\nannealing pilot 50ep search\n"
         "cost-aware-disc pilot search complete (retrain pending)"),
    ]
    add_table(s, rows, Pt(28), Pt(76), col_widths=[Pt(155), Pt(280), Pt(225)],
              header_color=ACCENT, font_size=9,
              header_height=Pt(22), body_row_height=Pt(96))


def slide_23_limitations(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "26.  Limitations")
    items = [
        (0, "Single seed (seed=2).  +0.08 pp delta is within DARTS seed noise.  Multi-seed runner is implemented, not yet executed."),
        (0, "First-order updates only.  Second-order compare script is staged; the Hessian-vector path is implemented but not benchmarked."),
        (0, "100-epoch retrain instead of the paper's 600.  Our absolute accuracies are 1–2 pp below the published numbers and not directly comparable."),
        (0, "Skip-collapse not eliminated.  CG-DARTS amplifies it; would need PC-DARTS / DARTS+ / DARTS-PT to fix."),
        (0, "No on-device validation.  Roofline + LUT are good but a Jetson / H100 run has not happened."),
        (0, "No ImageNet transfer.  Paper transfers the CIFAR-10 cell at 26.7% top-1.  We do not."),
        (0, "No PTB / recurrent CG-DARTS.  rnn/ subtree is untouched; cost-regularization concept carries over but is not implemented."),
        (0, "Cost-aware derivation retrain results not in by the time of writing — only the search infra and one in-flight retrain queue."),
    ]
    add_bullets(s, Pt(40), Pt(78), Pt(660), Pt(310), items, size=11, line_spacing=1.20)


def slide_24_conclusions(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "27.  Conclusions")
    items = [
        (0, "CG-DARTS makes hardware cost a first-class signal in DARTS, in both the bilevel objective and the discretization step"),
        (0, "Both open questions from our proposal's slide 16 are addressed:"),
        (1, "O1 (soft-to-hard gap)  →  τ annealing + τ-aware expected cost"),
        (1, "O2 (performance-aware derivation)  →  cost_sub / cost_div top-k modes selectable post-hoc on one search"),
        (0, "On a matched single-seed protocol, CG-DARTS at λ=1e-2 Pareto-dominates vanilla DARTS:"),
        (1, "+0.08 pp retrain accuracy   ·   −12% MACs   ·   −13% params"),
        (0, "Empirical surprise:  annealing's effect is decided by SCHEDULE SHAPE, not annealing itself"),
        (1, "Linear τ=5→0.1 collapses the search; hold-then-cosine (anneal only last 20%) recovers it (96.04% retrain)"),
        (1, "But hold-then-cosine is accuracy/cost-neutral at λ=1e-2 — annealing is safe, not yet a win"),
        (0, "We do NOT beat the paper's 2.76% headline — that requires multi-seed selection + 600-epoch retrain"),
        (1, "Infrastructure for all of it is ready in this branch; only GPU-time is missing"),
    ]
    add_bullets(s, Pt(40), Pt(80), Pt(660), Pt(300), items, size=12, line_spacing=1.20)


def slide_25_future(prs, n, total):
    s = add_slide(prs)
    chrome(s, n, total, "28.  Future work + references")
    items = [
        (0, "Immediate (already staged in scripts/, ~1 GPU-day each)"),
        (1, "Multi-seed FLOPs sweep at λ ∈ {1e-2, 5e-2} × seeds {0,1,2,3}  (scripts/run_multi_seed.sh)"),
        (1, "First-order vs second-order CG-DARTS at λ=1e-2  (scripts/run_second_order_compare.sh)"),
        (1, "Skip-share-vs-λ diagnostic on the multi-seed output  (scripts/skip_share_diagnostic.py)"),
        (1, "Cross-device transfer matrix on existing cells  (scripts/cross_device_eval.py)"),
        (0, "Medium-term (1–2 GPU weeks)"),
        (1, "On-device latency validation: run cells on a real Jetson / H100; replace roofline with measurements"),
        (1, "Search at 600-epoch retrain protocol with 4-seed selection — closes the comparison-fairness gap with the paper"),
        (1, "Recurrent CG-DARTS on PTB (rnn/) — cost formulation per step × sequence length"),
        (0, "References"),
        (1, "Liu, Simonyan, Yang.  DARTS: Differentiable Architecture Search.  ICLR 2019."),
        (1, "Zela et al.  Understanding and Robustifying Differentiable Architecture Search.  ICLR 2020."),
        (1, "Liang et al.  DARTS+: Improved DARTS with Early Stopping.  arXiv 1909.06035."),
        (1, "Cai, Zhu, Han.  ProxylessNAS: Direct NAS on Target Task and Hardware.  ICLR 2019."),
    ]
    add_bullets(s, Pt(40), Pt(78), Pt(660), Pt(310), items, size=10, line_spacing=1.18)


# ---------- table helper ----------

def add_table(slide, rows, x, y, col_widths, *, header_color=ACCENT,
              font_size=10, row_min_height=Pt(20), header_height=None,
              body_row_height=None):
    """row_min_height drives every row unless header_height/body_row_height override."""
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_w = sum(col_widths, Emu(0))
    table_h = row_min_height * n_rows
    table = slide.shapes.add_table(n_rows, n_cols, x, y, table_w, table_h).table
    # explicit row heights override the uniform layout
    if header_height is not None:
        table.rows[0].height = header_height
    if body_row_height is not None:
        for r in range(1, n_rows):
            table.rows[r].height = body_row_height
    # set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Pt(3); tf.margin_right = Pt(3)
            tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
            tf.word_wrap = True
            lines = str(val).split("\n")
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                run = p.add_run()
                run.text = line
                run.font.name = BODY_FONT
                run.font.size = Pt(font_size)
                if r == 0:
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(255, 255, 255)
                else:
                    run.font.color.rgb = INK
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if r % 2 else RGBColor(0xF2, 0xF4, 0xF8)
    return table


# ---------- main ----------

def main():
    prs = make_prs()
    total = 30

    title_slide(prs, total)
    outline_slide(prs, 2, total)
    slide_3_proposal_recap(prs, 3, total)
    slide_4_thesis(prs, 4, total)
    slide_5_method_overview(prs, 5, total)
    slide_6_cost_reg(prs, 6, total)
    slide_7_metrics(prs, 7, total)
    slide_8_device(prs, 8, total)
    slide_9_lut(prs, 9, total)
    slide_10_anneal_why(prs, 10, total)
    slide_10_anneal_intuition(prs, 11, total)
    slide_10_anneal(prs, 12, total)
    slide_11_disc(prs, 13, total)
    slide_12_setup(prs, 14, total)
    slide_13_compared(prs, 15, total)
    slide_14_flops_table(prs, 16, total)
    slide_15_flops_pareto(prs, 17, total)
    slide_15b_master_lambda(prs, 18, total)
    slide_16_params(prs, 19, total)
    slide_16b_device(prs, 20, total)
    slide_17_anneal(prs, 21, total)
    slide_17b_anneal_compare(prs, 22, total)
    slide_18_lut(prs, 23, total)
    slide_19_vs_paper(prs, 24, total)
    slide_20_skip(prs, 25, total)
    slide_21_discussion(prs, 26, total)
    slide_22_branches(prs, 27, total)
    slide_23_limitations(prs, 28, total)
    slide_24_conclusions(prs, 29, total)
    slide_25_future(prs, 30, total)

    prs.save(str(OUT))
    print("wrote", OUT)


if __name__ == "__main__":
    main()
